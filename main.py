# main.py - FastAPI Backend for RAG Workflow System
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Dict, Any, Optional
import os
import logging
from datetime import datetime
import uuid
import shutil
from pathlib import Path
import mimetypes

# Document processing
import PyMuPDF as fitz
from docx import Document
import pandas as pd
from pptx import Presentation

# Vector store and embeddings
import chromadb
import openai
import google.generativeai as genai
from sentence_transformers import SentenceTransformer

# Web search
import requests
from bs4 import BeautifulSoup
from urllib.parse import quote_plus

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# FastAPI app
app = FastAPI(title="RAG Workflow API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configs
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)
CHROMA_PERSIST_DIR = "./chroma_db"
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
chroma_client = chromadb.PersistentClient(path=CHROMA_PERSIST_DIR)

# Pydantic Models
class NodeData(BaseModel):
    label: str
    color: Optional[str] = None
    queryText: Optional[str] = None
    points: Optional[str] = None
    file: Optional[str] = None
    file_path: Optional[str] = None
    file_name: Optional[str] = None
    embeddingModel: Optional[str] = "text-embedding-3-large"
    uploadedFiles: Optional[List[Dict[str, Any]]] = []
    collection_name: Optional[str] = None
    chunk_size: Optional[int] = 1000
    chunk_overlap: Optional[int] = 200
    apiKey: Optional[str] = None
    llmModel: Optional[str] = "gpt-4o-mini"
    temperature: Optional[float] = 0.7
    prompt: Optional[str] = None
    webSearchTool: Optional[str] = "SerpAPI"
    maxTokens: Optional[int] = 1000
    topP: Optional[float] = 1.0
    frequencyPenalty: Optional[float] = 0
    presencePenalty: Optional[float] = 0
    outputText: Optional[str] = None
    title: Optional[str] = None
    allowFollowUp: Optional[bool] = False

class WorkflowNode(BaseModel):
    id: str
    type: str
    position: Dict[str, float]
    data: NodeData

class WorkflowEdge(BaseModel):
    id: str
    source: str
    target: str

class WorkflowRequest(BaseModel):
    nodes: List[WorkflowNode]
    edges: List[WorkflowEdge]
    query: str

# Document Processor
class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        try:
            doc = fitz.open(file_path)
            text = "".join(page.get_text() for page in doc)
            doc.close()
            return text
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return ""

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        try:
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return ""

    @staticmethod
    def extract_text_from_excel(file_path: str) -> str:
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            text = ""
            for sheet, df in df_dict.items():
                text += f"Sheet: {sheet}\n" + df.to_string(index=False) + "\n\n"
            return text
        except Exception as e:
            logger.error(f"Excel extraction error: {e}")
            return ""

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return ""

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"TXT extraction error: {e}")
            return ""

    @classmethod
    def process_file(cls, file_path: str, content_type: str) -> str:
        content_type = content_type.lower()
        if 'pdf' in content_type:
            return cls.extract_text_from_pdf(file_path)
        elif 'word' in content_type or 'docx' in content_type:
            return cls.extract_text_from_docx(file_path)
        elif 'excel' in content_type or 'spreadsheet' in content_type:
            return cls.extract_text_from_excel(file_path)
        elif 'presentation' in content_type or 'ppt' in content_type:
            return cls.extract_text_from_pptx(file_path)
        elif content_type.startswith('text/'):
            return cls.extract_text_from_txt(file_path)
        else:
            logger.warning(f"Unknown file type: {content_type}, treating as text")
            return cls.extract_text_from_txt(file_path)

# Embedding Service
class EmbeddingService:
    def __init__(self):
        self.sentence_transformer = None

    def get_openai_embeddings(self, texts: List[str], model: str, api_key: str) -> List[List[float]]:
        try:
            client = openai.OpenAI(api_key=api_key)
            embeddings = []
            for i in range(0, len(texts), 50):
                batch = texts[i:i+50]
                response = client.embeddings.create(model=model, input=batch)
                embeddings.extend([item.embedding for item in response.data])
            return embeddings
        except Exception as e:
            logger.error(f"OpenAI embedding error: {e}")
            raise HTTPException(status_code=500, detail=str(e))

    def get_gemini_embeddings(self, texts: List[str], api_key: str) -> List[List[float]]:
        try:
            genai.configure(api_key=api_key)
            embeddings = []
            for text in texts:
                result = genai.embed_content(model="models/embedding-001",
                                             content=text,
                                             task_type="retrieval_document")
                embeddings.append(result['embedding'])
            return embeddings
        except Exception as e:
            logger.error(f"Gemini embedding error: {e}")
            raise HTTPException(status_code=500, detail=str(e))

# Vector Store
class VectorStore:
    def __init__(self, chroma_client):
        self.client = chroma_client
        self.embedding_service = EmbeddingService()

    def chunk_text(self, text: str, chunk_size=1000, chunk_overlap=200) -> List[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - chunk_overlap):
            chunks.append(" ".join(words[i:i+chunk_size]))
        return chunks

    def store_document(self, file_path: str, text: str, embedding_model: str, api_key: str, collection_name: str, chunk_size=1000, chunk_overlap=200):
        try:
            collection = self.client.get_or_create_collection(name=collection_name)
            chunks = self.chunk_text(text, chunk_size, chunk_overlap)
            if embedding_model.startswith('text-embedding'):
                embeddings = self.embedding_service.get_openai_embeddings(chunks, embedding_model, api_key)
            elif embedding_model == 'gemini':
                embeddings = self.embedding_service.get_gemini_embeddings(chunks, api_key)
            else:
                if not self.embedding_service.sentence_transformer:
                    self.embedding_service.sentence_transformer = SentenceTransformer('all-MiniLM-L6-v2')
                embeddings = self.embedding_service.sentence_transformer.encode(chunks).tolist()
            ids = [f"{file_path}_{i}" for i in range(len(chunks))]
            metadatas = [{"file_path": file_path, "chunk_index": i} for i in range(len(chunks))]
            collection.add(embeddings=embeddings, documents=chunks, metadatas=metadatas, ids=ids)
            return True
        except Exception as e:
            logger.error(f"Vector store error: {e}")
            raise HTTPException(status_code=500, detail=str(e))

    def query_similar(self, query: str, embedding_model: str, api_key: str, collection_name: str, n_results=5) -> List[str]:
        try:
            collection = self.client.get_collection(collection_name)
            if embedding_model.startswith('text-embedding'):
                query_embedding = self.embedding_service.get_openai_embeddings([query], embedding_model, api_key)[0]
            elif embedding_model == 'gemini':
                query_embedding = self.embedding_service.get_gemini_embeddings([query], api_key)[0]
            else:
                if not self.embedding_service.sentence_transformer:
                    self.embedding_service.sentence_transformer = SentenceTransformer('all-MiniLM-L6-v2')
                query_embedding = self.embedding_service.sentence_transformer.encode([query])[0].tolist()
            results = collection.query(query_embeddings=[query_embedding], n_results=n_results)
            return results['documents'][0] if results['documents'] else []
        except Exception as e:
            logger.error(f"Query error: {e}")
            return []

# LLM and Web Search Services are unchanged

vector_store = VectorStore(chroma_client)

# Remaining endpoints (/upload_file, /run_workflow, /health) remain same as your version
@app.post("/upload_file")
async def upload_file(file: UploadFile = File(...)):
    """Handle file upload from frontend"""
    try:
        # Read file content to check size
        file_contents = await file.read()
        file_size = len(file_contents)
        
        if file_size > MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail="File too large")
        
        # Generate unique filename
        file_id = str(uuid.uuid4())
        file_extension = Path(file.filename).suffix
        unique_filename = f"{file_id}{file_extension}"
        file_path = UPLOAD_DIR / unique_filename
        
        # Save file to disk
        with open(file_path, "wb") as buffer:
            buffer.write(file_contents)
        
        # Extract preview text
        content_type = file.content_type or mimetypes.guess_type(file.filename)[0] or ""
        preview_text = DocumentProcessor.process_file(str(file_path), content_type)
        preview = preview_text[:200] if preview_text else ""
        
        logger.info(f"Uploaded file: {file.filename} -> {file_path}")
        
        return {
            "filename": file.filename,
            "file_path": str(file_path),
            "size": file_size,
            "content_type": content_type,
            "preview": preview
        }
        
    except Exception as e:
        logger.error(f"File upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")



















